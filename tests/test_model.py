"""Tests for the model layer — snapshot, index, refs."""

import pytest

from pptx import Presentation

from fcp_slides.model.snapshot import SlidesModel, SnapshotEvent
from fcp_slides.model.index import SlideIndex
from fcp_slides.model.refs import ShapeRef


class TestSlidesModel:
    def test_create_default(self):
        model = SlidesModel()
        assert model.title == "Untitled"
        assert model.prs is not None
        assert model.file_path is None

    def test_create_with_title(self):
        model = SlidesModel(title="Test Deck")
        assert model.title == "Test Deck"

    def test_snapshot_roundtrip(self):
        model = SlidesModel(title="Test")
        # Add a slide
        layout = model.prs.slide_layouts[0]
        model.prs.slides.add_slide(layout)
        assert len(model.prs.slides) == 1

        # Snapshot
        data = model.snapshot()
        assert isinstance(data, bytes)
        assert len(data) > 0

        # Restore
        model.restore(data)
        assert len(model.prs.slides) == 1

    def test_snapshot_undo(self):
        model = SlidesModel(title="Test")
        layout = model.prs.slide_layouts[0]

        # State 1: no slides
        before = model.snapshot()

        # State 2: one slide
        model.prs.slides.add_slide(layout)
        assert len(model.prs.slides) == 1

        # Undo
        model.restore(before)
        assert len(model.prs.slides) == 0


class TestSnapshotEvent:
    def test_defaults(self):
        event = SnapshotEvent()
        assert event.type == "snapshot"
        assert event.before == b""
        assert event.after == b""
        assert event.summary == ""


class TestSlideIndex:
    def setup_method(self):
        self.index = SlideIndex()

    def test_initial_state(self):
        assert self.index.active_slide == 0
        assert self.index._slide_labels == {}
        assert self.index._shape_labels == {}

    def test_add_slide_label(self):
        self.index.add_slide_label("intro", 0)
        assert "intro" in self.index._slide_labels

    def test_remove_slide_label(self):
        self.index.add_slide_label("intro", 0)
        self.index.remove_slide_label("intro")
        assert "intro" not in self.index._slide_labels

    def test_add_shape_label(self):
        ref = ShapeRef(label="title1", slide_idx=0, shape_id=1, shape_type="text_box")
        self.index.add_shape_label("title1", ref)
        assert self.index.resolve_shape("title1") is not None

    def test_resolve_slide_idx_by_number(self):
        model = SlidesModel()
        layout = model.prs.slide_layouts[0]
        model.prs.slides.add_slide(layout)
        model.prs.slides.add_slide(layout)
        self.index.rebuild(model)

        assert self.index.resolve_slide_idx("1", model) == 0
        assert self.index.resolve_slide_idx("2", model) == 1
        assert self.index.resolve_slide_idx("3", model) is None

    def test_resolve_slide_idx_by_label(self):
        model = SlidesModel()
        layout = model.prs.slide_layouts[0]
        model.prs.slides.add_slide(layout)
        self.index.rebuild(model)

        assert self.index.resolve_slide_idx("s1", model) == 0

    def test_resolve_slide_idx_active(self):
        model = SlidesModel()
        layout = model.prs.slide_layouts[0]
        model.prs.slides.add_slide(layout)
        self.index.rebuild(model)
        self.index.active_slide = 0

        assert self.index.resolve_slide_idx("active", model) == 0

    def test_resolve_slide_idx_last(self):
        model = SlidesModel()
        layout = model.prs.slide_layouts[0]
        model.prs.slides.add_slide(layout)
        model.prs.slides.add_slide(layout)
        self.index.rebuild(model)

        assert self.index.resolve_slide_idx("last", model) == 1

    def test_rebuild(self):
        model = SlidesModel()
        layout = model.prs.slide_layouts[0]
        model.prs.slides.add_slide(layout)
        model.prs.slides.add_slide(layout)

        self.index.rebuild(model)

        assert "s1" in self.index._slide_labels
        assert "s2" in self.index._slide_labels

    def test_clear(self):
        self.index.add_slide_label("test", 0)
        self.index.clear()
        assert self.index._slide_labels == {}
        assert self.index._shape_labels == {}
        assert self.index.active_slide == 0


class TestShapeRef:
    def test_creation(self):
        ref = ShapeRef(
            label="my_box",
            slide_idx=0,
            shape_id=42,
            shape_type="rectangle",
        )
        assert ref.label == "my_box"
        assert ref.slide_idx == 0
        assert ref.shape_id == 42
        assert ref.placeholder_idx is None
