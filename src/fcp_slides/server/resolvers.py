"""Operation context and reference resolution for slides verbs.

Provides SlidesOpContext and helpers for resolving slide/shape references
and extracting position/size parameters.
"""

from __future__ import annotations

from dataclasses import dataclass

from pptx import Presentation
from pptx.slide import Slide

from fcp_slides.lib.units import parse_length
from fcp_slides.model.index import SlideIndex
from fcp_slides.model.refs import ShapeRef
from fcp_slides.model.snapshot import SlidesModel


@dataclass
class SlidesOpContext:
    """Context passed to every verb handler."""

    prs: Presentation
    index: SlideIndex
    model: SlidesModel

    @property
    def active_slide(self) -> Slide | None:
        """The currently active slide."""
        slides = self.prs.slides
        idx = self.index.active_slide
        if 0 <= idx < len(slides):
            return slides[idx]
        return None

    @property
    def active_slide_idx(self) -> int:
        return self.index.active_slide

    @property
    def slide_count(self) -> int:
        return len(self.prs.slides)


def resolve_slide(ref: str, ctx: SlidesOpContext) -> tuple[Slide, int] | None:
    """Resolve a slide reference to (Slide, 0-based index).

    Returns None if not found.
    """
    idx = ctx.index.resolve_slide_idx(ref, ctx.model)
    if idx is None:
        return None
    slides = ctx.prs.slides
    if 0 <= idx < len(slides):
        return slides[idx], idx
    return None


def resolve_shape_on_slide(
    ref: str, slide: Slide, slide_idx: int, ctx: SlidesOpContext
) -> "BaseShape | None":
    """Resolve a shape reference on a specific slide.

    Tries:
      1. Index label lookup
      2. Shape name match
      3. 1-based shape index
    """
    # Try index label
    shape_ref = ctx.index.resolve_shape(ref, slide_idx)
    if shape_ref is not None:
        for shape in slide.shapes:
            if shape.shape_id == shape_ref.shape_id:
                return shape

    # Try shape name match
    ref_lower = ref.lower()
    for shape in slide.shapes:
        if shape.name.lower() == ref_lower:
            return shape

    # Try 1-based index
    try:
        n = int(ref)
        shapes_list = list(slide.shapes)
        if 1 <= n <= len(shapes_list):
            return shapes_list[n - 1]
    except ValueError:
        pass

    return None


def extract_position(params: dict[str, str]) -> dict[str, int]:
    """Extract position/size EMU values from params.

    Recognized params: x, y, w (width), h (height), cx (alias for w), cy (alias for h).
    Returns dict with keys: left, top, width, height (only those present).
    """
    result: dict[str, int] = {}

    if "x" in params:
        result["left"] = parse_length(params["x"])
    if "y" in params:
        result["top"] = parse_length(params["y"])

    w = params.get("w") or params.get("cx") or params.get("width")
    if w:
        result["width"] = parse_length(w)

    h = params.get("h") or params.get("cy") or params.get("height")
    if h:
        result["height"] = parse_length(h)

    return result


def require_active_slide(ctx: SlidesOpContext) -> tuple[Slide, int] | str:
    """Get the active slide or return an error message."""
    slide = ctx.active_slide
    if slide is None:
        return "No active slide. Use 'slide add' first."
    return slide, ctx.active_slide_idx
