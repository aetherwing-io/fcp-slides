"""Chart operation handlers — add, data, series, axis, remove."""

from __future__ import annotations

from pptx.util import Inches, Emu
from pptx.chart.data import CategoryChartData

from fcp_core import OpResult, ParsedOp

from fcp_slides.lib.chart_types import resolve_chart_type, list_chart_types
from fcp_slides.model.refs import ShapeRef
from fcp_slides.server.resolvers import (
    SlidesOpContext,
    extract_position,
    require_active_slide,
    resolve_shape_on_slide,
)


_DEFAULT_CHART_LEFT = Inches(1)
_DEFAULT_CHART_TOP = Inches(2)
_DEFAULT_CHART_WIDTH = Inches(8)
_DEFAULT_CHART_HEIGHT = Inches(4.5)


def op_chart(op: ParsedOp, ctx: SlidesOpContext) -> OpResult:
    """Dispatch chart sub-commands."""
    if not op.positionals:
        return OpResult(success=False, message="Usage: chart add|data|series|axis|remove ...")

    action = op.positionals[0].lower()
    rest = op.positionals[1:]

    dispatch = {
        "add": _chart_add,
        "data": _chart_data,
        "series": _chart_series,
        "axis": _chart_axis,
        "remove": _chart_remove,
    }

    handler = dispatch.get(action)
    if handler is None:
        return OpResult(
            success=False,
            message=f"Unknown chart action: {action!r}. Use: add, data, series, axis, remove",
        )

    return handler(rest, op.params, ctx)


def _chart_add(
    args: list[str], params: dict[str, str], ctx: SlidesOpContext
) -> OpResult:
    """Add a chart to the active slide.

    Usage: chart add TYPE [label:NAME] [title:TITLE] [x:POS] [y:POS] [w:SIZE] [h:SIZE]
    """
    result = require_active_slide(ctx)
    if isinstance(result, str):
        return OpResult(success=False, message=result)
    slide, slide_idx = result

    if not args:
        return OpResult(
            success=False,
            message=f"Usage: chart add TYPE [label:NAME]\nTypes: {', '.join(list_chart_types()[:10])}...",
        )

    chart_type_name = args[0]
    chart_type = resolve_chart_type(chart_type_name)
    if chart_type is None:
        return OpResult(
            success=False,
            message=f"Unknown chart type: {chart_type_name!r}. Use: {', '.join(list_chart_types())}",
        )

    pos = extract_position(params)
    left = pos.get("left", _DEFAULT_CHART_LEFT)
    top = pos.get("top", _DEFAULT_CHART_TOP)
    width = pos.get("width", _DEFAULT_CHART_WIDTH)
    height = pos.get("height", _DEFAULT_CHART_HEIGHT)

    # Create chart with empty placeholder data
    chart_data = CategoryChartData()
    chart_data.categories = ["Category 1"]
    chart_data.add_series("Series 1", (0,))

    shape = slide.shapes.add_chart(chart_type, left, top, width, height, chart_data)
    chart = shape.chart

    # Set title if provided
    title = params.get("title", "")
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title

    label = params.get("label", "")
    if label:
        ref = ShapeRef(
            label=label,
            slide_idx=slide_idx,
            shape_id=shape.shape_id,
            shape_type="chart",
        )
        ctx.index.add_shape_label(label, ref)
    else:
        ctx.index.rebuild(ctx.model)

    return OpResult(
        success=True,
        message=f"Chart '{chart_type_name}' added" + (f" label:{label}" if label else ""),
        prefix="+",
    )


def _resolve_chart_shape(ref: str, ctx: SlidesOpContext):
    """Find a chart shape by label on the active slide."""
    active = require_active_slide(ctx)
    if isinstance(active, str):
        return None, None, active
    slide, slide_idx = active

    shape = resolve_shape_on_slide(ref, slide, slide_idx, ctx)
    if shape is None:
        return None, None, f"Shape not found: {ref!r}"

    if not shape.has_chart:
        return None, None, f"Shape '{ref}' is not a chart"

    return shape, shape.chart, None


def _chart_data(
    args: list[str], params: dict[str, str], ctx: SlidesOpContext
) -> OpResult:
    """Set chart data (replace all data).

    Usage: chart data CHART_REF categories:"Cat1,Cat2" series:"S1" values:"1,2"
    """
    if not args:
        return OpResult(success=False, message="Usage: chart data CHART_REF categories:\"...\" series:\"...\" values:\"...\"")

    chart_ref = args[0]
    shape, chart, err = _resolve_chart_shape(chart_ref, ctx)
    if err:
        return OpResult(success=False, message=err)

    categories_str = params.get("categories", "")
    series_name = params.get("series", "Series 1")
    values_str = params.get("values", "")

    if not categories_str or not values_str:
        return OpResult(success=False, message="Both categories: and values: are required")

    categories = [c.strip() for c in categories_str.split(",")]
    values = _parse_values(values_str)

    if len(values) != len(categories):
        return OpResult(
            success=False,
            message=f"Mismatched counts: {len(categories)} categories, {len(values)} values",
        )

    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series(series_name, values)

    chart.replace_data(chart_data)

    return OpResult(
        success=True,
        message=f"Chart data set: {len(categories)} categories, series '{series_name}'",
        prefix="*",
    )


def _chart_series(
    args: list[str], params: dict[str, str], ctx: SlidesOpContext
) -> OpResult:
    """Add a data series to an existing chart.

    Usage: chart series CHART_REF series:"S2" values:"3,4,5"
    """
    if not args:
        return OpResult(success=False, message='Usage: chart series CHART_REF series:"NAME" values:"1,2,3"')

    chart_ref = args[0]
    shape, chart, err = _resolve_chart_shape(chart_ref, ctx)
    if err:
        return OpResult(success=False, message=err)

    series_name = params.get("series", "")
    values_str = params.get("values", "")

    if not series_name or not values_str:
        return OpResult(success=False, message="Both series: and values: are required")

    values = _parse_values(values_str)

    # Get existing data and add new series
    plot = chart.plots[0]
    existing_categories = [str(pt.label) for pt in plot.categories]

    chart_data = CategoryChartData()
    chart_data.categories = existing_categories

    # Re-add existing series
    for series in plot.series:
        series_vals = [series.values[i] if i < len(series.values) else 0 for i in range(len(existing_categories))]
        chart_data.add_series(str(series.name) if series.name else "Series", series_vals)

    # Add new series
    chart_data.add_series(series_name, values)
    chart.replace_data(chart_data)

    return OpResult(
        success=True,
        message=f"Series '{series_name}' added ({len(values)} values)",
        prefix="+",
    )


def _chart_axis(
    args: list[str], params: dict[str, str], ctx: SlidesOpContext
) -> OpResult:
    """Configure chart axes.

    Usage: chart axis CHART_REF x|y|value|category [title:"TITLE"] [min:N] [max:N]
    """
    if len(args) < 2:
        return OpResult(success=False, message="Usage: chart axis CHART_REF x|y [title:\"TITLE\"]")

    chart_ref = args[0]
    axis_name = args[1].lower()
    shape, chart, err = _resolve_chart_shape(chart_ref, ctx)
    if err:
        return OpResult(success=False, message=err)

    # Map axis names
    if axis_name in ("x", "category"):
        axis = chart.category_axis
    elif axis_name in ("y", "value"):
        axis = chart.value_axis
    else:
        return OpResult(success=False, message=f"Unknown axis: {axis_name!r}. Use: x, y, category, value")

    changes: list[str] = []

    if "title" in params:
        axis.has_title = True
        axis.axis_title.text_frame.text = params["title"]
        changes.append(f"title=\"{params['title']}\"")

    if "min" in params:
        try:
            axis.minimum_scale = float(params["min"])
            changes.append(f"min={params['min']}")
        except (ValueError, AttributeError):
            pass

    if "max" in params:
        try:
            axis.maximum_scale = float(params["max"])
            changes.append(f"max={params['max']}")
        except (ValueError, AttributeError):
            pass

    if not changes:
        return OpResult(success=False, message="No axis properties specified. Use: title, min, max")

    return OpResult(
        success=True,
        message=f"Axis '{axis_name}' updated: {', '.join(changes)}",
        prefix="*",
    )


def _chart_remove(
    args: list[str], params: dict[str, str], ctx: SlidesOpContext
) -> OpResult:
    """Remove a chart (shape) from the slide."""
    if not args:
        return OpResult(success=False, message="Usage: chart remove CHART_REF")

    chart_ref = args[0]
    active = require_active_slide(ctx)
    if isinstance(active, str):
        return OpResult(success=False, message=active)
    slide, slide_idx = active

    shape = resolve_shape_on_slide(chart_ref, slide, slide_idx, ctx)
    if shape is None:
        return OpResult(success=False, message=f"Shape not found: {chart_ref!r}")

    sp = shape._element
    sp.getparent().remove(sp)
    ctx.index.rebuild(ctx.model)

    return OpResult(success=True, message=f"Chart '{chart_ref}' removed", prefix="-")


def _parse_values(values_str: str) -> list[float]:
    """Parse comma-separated values, handling suffixes like M, K."""
    values: list[float] = []
    for v in values_str.split(","):
        v = v.strip()
        if not v:
            values.append(0)
            continue

        # Handle suffixes
        multiplier = 1.0
        if v.upper().endswith("M"):
            v = v[:-1]
            multiplier = 1_000_000
        elif v.upper().endswith("K"):
            v = v[:-1]
            multiplier = 1_000

        # Strip $ and other currency symbols
        v = v.lstrip("$").strip()

        try:
            values.append(float(v) * multiplier)
        except ValueError:
            values.append(0)

    return values


HANDLERS: dict[str, callable] = {
    "chart": op_chart,
}
