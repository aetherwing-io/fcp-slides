"""Read-only query handlers for slides inspection.

Dispatches query strings to handlers that return formatted
information about the presentation state.
"""

from __future__ import annotations

from pptx.util import Emu

from fcp_slides.lib.units import format_length
from fcp_slides.model.index import SlideIndex, _shape_type_name
from fcp_slides.model.snapshot import SlidesModel


def dispatch_query(query: str, model: SlidesModel, index: SlideIndex) -> str:
    """Route a query string to the appropriate handler."""
    parts = query.strip().split(None, 1)
    if not parts:
        return "! Empty query"

    cmd = parts[0].lower()
    args = parts[1] if len(parts) > 1 else ""

    handler = QUERY_HANDLERS.get(cmd)
    if handler is None:
        available = ", ".join(sorted(QUERY_HANDLERS.keys()))
        return f"! Unknown query: {cmd!r}. Available: {available}"

    return handler(args, model, index)


def _query_plan(args: str, model: SlidesModel, index: SlideIndex) -> str:
    """Overview of the presentation structure."""
    prs = model.prs
    slides = prs.slides
    lines: list[str] = []
    lines.append(f"Presentation: {model.title}")
    if model.file_path:
        lines.append(f"File: {model.file_path}")

    w = prs.slide_width
    h = prs.slide_height
    if w and h:
        lines.append(f"Size: {format_length(w)} x {format_length(h)}")

    lines.append(f"Slides: {len(slides)}")
    lines.append(f"Active: slide {index.active_slide + 1}")
    lines.append("")

    for i, slide in enumerate(slides):
        marker = " *" if i == index.active_slide else ""
        layout_name = slide.slide_layout.name if slide.slide_layout else "?"
        lines.append(f"  Slide {i + 1} [{layout_name}]{marker}")

        for shape in slide.shapes:
            type_name = _shape_type_name(shape)
            pos = ""
            if shape.left is not None and shape.top is not None:
                pos = f" at ({format_length(shape.left)}, {format_length(shape.top)})"
            size = ""
            if shape.width is not None and shape.height is not None:
                size = f" {format_length(shape.width)}x{format_length(shape.height)}"

            text_preview = ""
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    text_preview = f' "{text[:40]}"'

            ph = ""
            if shape.is_placeholder and shape.placeholder_format:
                ph = f" [ph:{shape.placeholder_format.idx}]"

            lines.append(f"    {type_name}{ph}: {shape.name}{pos}{size}{text_preview}")

    return "\n".join(lines)


def _query_status(args: str, model: SlidesModel, index: SlideIndex) -> str:
    """Quick status summary."""
    prs = model.prs
    slide_count = len(prs.slides)
    shape_count = sum(len(list(s.shapes)) for s in prs.slides)
    lines = [
        f"Title: {model.title}",
        f"File: {model.file_path or '(unsaved)'}",
        f"Slides: {slide_count}",
        f"Shapes: {shape_count}",
        f"Active: slide {index.active_slide + 1}",
    ]
    return "\n".join(lines)


def _query_describe(args: str, model: SlidesModel, index: SlideIndex) -> str:
    """Describe a specific slide or shape in detail."""
    if not args:
        return "! Usage: describe SLIDE_REF [SHAPE_REF]"

    parts = args.strip().split(None, 1)
    slide_ref = parts[0]
    shape_ref = parts[1] if len(parts) > 1 else None

    slide_idx = index.resolve_slide_idx(slide_ref, model)
    if slide_idx is None:
        return f"! Slide not found: {slide_ref!r}"

    slide = model.prs.slides[slide_idx]

    if shape_ref:
        # Describe a specific shape
        from fcp_slides.server.resolvers import resolve_shape_on_slide, SlidesOpContext

        ctx = SlidesOpContext(prs=model.prs, index=index, model=model)
        shape = resolve_shape_on_slide(shape_ref, slide, slide_idx, ctx)
        if shape is None:
            return f"! Shape not found: {shape_ref!r} on slide {slide_idx + 1}"

        return _describe_shape(shape, slide_idx)

    # Describe the slide
    return _describe_slide(slide, slide_idx, index)


def _describe_slide(slide, slide_idx: int, index: SlideIndex) -> str:
    """Detailed description of a slide."""
    lines: list[str] = []
    layout_name = slide.slide_layout.name if slide.slide_layout else "?"
    lines.append(f"Slide {slide_idx + 1} — Layout: {layout_name}")

    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
        notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            lines.append(f"Notes: {notes[:100]}")

    shapes = list(slide.shapes)
    lines.append(f"Shapes: {len(shapes)}")
    lines.append("")

    for shape in shapes:
        type_name = _shape_type_name(shape)
        lines.append(f"  [{type_name}] {shape.name}")
        if shape.left is not None:
            lines.append(f"    Position: ({format_length(shape.left)}, {format_length(shape.top)})")
        if shape.width is not None:
            lines.append(f"    Size: {format_length(shape.width)} x {format_length(shape.height)}")
        if shape.is_placeholder and shape.placeholder_format:
            lines.append(f"    Placeholder: idx={shape.placeholder_format.idx}")
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                lines.append(f"    Text: {text[:80]}")
        if shape.has_table:
            tbl = shape.table
            lines.append(f"    Table: {len(tbl.rows)} rows x {len(tbl.columns)} cols")
        if shape.has_chart:
            lines.append(f"    Chart: {shape.chart.chart_type}")

    return "\n".join(lines)


def _describe_shape(shape, slide_idx: int) -> str:
    """Detailed description of a single shape."""
    type_name = _shape_type_name(shape)
    lines: list[str] = [
        f"Shape: {shape.name} [{type_name}] on slide {slide_idx + 1}",
    ]

    if shape.left is not None:
        lines.append(f"Position: ({format_length(shape.left)}, {format_length(shape.top)})")
    if shape.width is not None:
        lines.append(f"Size: {format_length(shape.width)} x {format_length(shape.height)}")
    if hasattr(shape, "rotation") and shape.rotation:
        lines.append(f"Rotation: {shape.rotation}")

    if shape.is_placeholder and shape.placeholder_format:
        lines.append(f"Placeholder: idx={shape.placeholder_format.idx}")

    if shape.has_text_frame:
        tf = shape.text_frame
        lines.append(f"Text paragraphs: {len(tf.paragraphs)}")
        for i, para in enumerate(tf.paragraphs):
            text = para.text.strip()
            if text:
                level = para.level or 0
                lines.append(f"  [{i}] L{level}: {text[:80]}")

    if shape.has_table:
        tbl = shape.table
        lines.append(f"Table: {len(tbl.rows)} rows x {len(tbl.columns)} cols")
        # Show first few rows
        for ri, row in enumerate(tbl.rows):
            if ri >= 5:
                lines.append(f"  ... +{len(tbl.rows) - 5} more rows")
                break
            cells = [tbl.cell(ri, ci).text[:20] for ci in range(len(tbl.columns))]
            lines.append(f"  Row {ri}: {' | '.join(cells)}")

    return "\n".join(lines)


def _query_list(args: str, model: SlidesModel, index: SlideIndex) -> str:
    """List slides, shapes, layouts, or labels."""
    what = args.strip().lower() if args else "slides"

    if what == "slides":
        lines: list[str] = []
        for i, slide in enumerate(model.prs.slides):
            marker = " *" if i == index.active_slide else ""
            layout = slide.slide_layout.name if slide.slide_layout else "?"
            shapes = len(list(slide.shapes))
            lines.append(f"  {i + 1}. [{layout}] {shapes} shapes{marker}")
        return "\n".join(lines) if lines else "No slides"

    if what == "layouts":
        from fcp_slides.lib.layout_names import list_layouts
        names = list_layouts(model.prs)
        return "\n".join(f"  {i}. {name}" for i, name in enumerate(names))

    if what == "shapes":
        slide = model.prs.slides[index.active_slide] if model.prs.slides else None
        if slide is None:
            return "No active slide"
        lines = []
        for shape in slide.shapes:
            type_name = _shape_type_name(shape)
            text_hint = ""
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    text_hint = f' "{t[:30]}"'
            lines.append(f"  {shape.name} [{type_name}]{text_hint}")
        return "\n".join(lines) if lines else "No shapes on active slide"

    if what == "labels":
        lines = []
        for label, idx in sorted(index._slide_labels.items()):
            lines.append(f"  slide: {label} → slide {idx + 1}")
        for label, ref in sorted(index._shape_labels.items()):
            lines.append(f"  shape: {label} → slide {ref.slide_idx + 1}, {ref.shape_type}")
        return "\n".join(lines) if lines else "No labels registered"

    return f"! Unknown list target: {what!r}. Use: slides, shapes, layouts, labels"


def _query_find(args: str, model: SlidesModel, index: SlideIndex) -> str:
    """Search for text across all slides."""
    if not args:
        return "! Usage: find TEXT"

    needle = args.strip().lower()
    results: list[str] = []

    for i, slide in enumerate(model.prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                if needle in text.lower():
                    preview = text.strip()[:60]
                    results.append(f"  Slide {i + 1}, {shape.name}: \"{preview}\"")

    if not results:
        return f"No matches for {args!r}"

    if len(results) > 50:
        results = results[:50]
        results.append(f"  ... truncated (50/{len(results)})")

    return "\n".join(results)


QUERY_HANDLERS: dict[str, callable] = {
    "plan": _query_plan,
    "map": _query_plan,
    "status": _query_status,
    "describe": _query_describe,
    "list": _query_list,
    "find": _query_find,
}
