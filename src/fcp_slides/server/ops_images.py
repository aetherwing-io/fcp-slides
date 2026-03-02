"""Image operation handlers — add, placeholder, remove."""

from __future__ import annotations

import os

from pptx.util import Inches

from fcp_core import OpResult, ParsedOp

from fcp_slides.model.refs import ShapeRef
from fcp_slides.server.resolvers import (
    SlidesOpContext,
    extract_position,
    require_active_slide,
    resolve_shape_on_slide,
    resolve_slide,
)


_DEFAULT_IMAGE_LEFT = Inches(2)
_DEFAULT_IMAGE_TOP = Inches(2)


def op_image(op: ParsedOp, ctx: SlidesOpContext) -> OpResult:
    """Dispatch image sub-commands."""
    if not op.positionals:
        return OpResult(success=False, message="Usage: image add|placeholder|remove PATH|REF ...")

    action = op.positionals[0].lower()
    rest = op.positionals[1:]

    dispatch = {
        "add": _image_add,
        "placeholder": _image_placeholder,
        "remove": _image_remove,
    }

    handler = dispatch.get(action)
    if handler is None:
        return OpResult(success=False, message=f"Unknown image action: {action!r}. Use: add, placeholder, remove")

    return handler(rest, op.params, ctx)


def _image_add(
    args: list[str], params: dict[str, str], ctx: SlidesOpContext
) -> OpResult:
    """Add an image to the active slide.

    Usage: image add PATH [x:POS] [y:POS] [w:SIZE] [h:SIZE] [label:NAME]
    """
    result = require_active_slide(ctx)
    if isinstance(result, str):
        return OpResult(success=False, message=result)
    slide, slide_idx = result

    if not args:
        return OpResult(success=False, message="Usage: image add PATH [x:POS] [y:POS] [w:SIZE] [h:SIZE]")

    image_path = args[0]
    if not os.path.isfile(image_path):
        return OpResult(success=False, message=f"Image file not found: {image_path!r}")

    pos = extract_position(params)
    left = pos.get("left", _DEFAULT_IMAGE_LEFT)
    top = pos.get("top", _DEFAULT_IMAGE_TOP)
    width = pos.get("width")
    height = pos.get("height")

    if width and height:
        shape = slide.shapes.add_picture(image_path, left, top, width, height)
    elif width:
        shape = slide.shapes.add_picture(image_path, left, top, width=width)
    elif height:
        shape = slide.shapes.add_picture(image_path, left, top, height=height)
    else:
        shape = slide.shapes.add_picture(image_path, left, top)

    label = params.get("label", "")
    if label:
        ref = ShapeRef(
            label=label,
            slide_idx=slide_idx,
            shape_id=shape.shape_id,
            shape_type="picture",
        )
        ctx.index.add_shape_label(label, ref)
    else:
        ctx.index.rebuild(ctx.model)

    fname = os.path.basename(image_path)
    return OpResult(
        success=True,
        message=f"Image '{fname}' added" + (f" label:{label}" if label else ""),
        prefix="+",
    )


def _image_placeholder(
    args: list[str], params: dict[str, str], ctx: SlidesOpContext
) -> OpResult:
    """Insert an image into a picture placeholder.

    Usage: image placeholder PATH [on:SLIDE]
    """
    if not args:
        return OpResult(success=False, message="Usage: image placeholder PATH [on:SLIDE]")

    image_path = args[0]
    if not os.path.isfile(image_path):
        return OpResult(success=False, message=f"Image file not found: {image_path!r}")

    on_ref = params.get("on")
    if on_ref:
        result = resolve_slide(on_ref, ctx)
        if result is None:
            return OpResult(success=False, message=f"Slide not found: {on_ref!r}")
        slide, slide_idx = result
    else:
        active = require_active_slide(ctx)
        if isinstance(active, str):
            return OpResult(success=False, message=active)
        slide, slide_idx = active

    # Find picture placeholder (idx 18 is common for picture placeholders)
    pic_ph = None
    for ph in slide.placeholders:
        # Picture placeholder types: 18 (picture), or check the XML
        ph_type = ph.placeholder_format.type
        if ph_type is not None and hasattr(ph_type, 'real') and ph_type.real == 18:
            pic_ph = ph
            break

    if pic_ph is None:
        # Try any placeholder that looks like a picture placeholder
        for ph in slide.placeholders:
            if "picture" in ph.name.lower() or "image" in ph.name.lower():
                pic_ph = ph
                break

    if pic_ph is None:
        return OpResult(success=False, message=f"No picture placeholder found on slide {slide_idx + 1}")

    pic_ph.insert_picture(image_path)
    fname = os.path.basename(image_path)
    return OpResult(success=True, message=f"Image '{fname}' inserted into placeholder", prefix="*")


def _image_remove(
    args: list[str], params: dict[str, str], ctx: SlidesOpContext
) -> OpResult:
    """Remove an image (shape) from the slide."""
    if not args:
        return OpResult(success=False, message="Usage: image remove SHAPE_REF")

    ref = args[0]
    active = require_active_slide(ctx)
    if isinstance(active, str):
        return OpResult(success=False, message=active)
    slide, slide_idx = active

    shape = resolve_shape_on_slide(ref, slide, slide_idx, ctx)
    if shape is None:
        return OpResult(success=False, message=f"Shape not found: {ref!r}")

    sp = shape._element
    sp.getparent().remove(sp)
    ctx.index.rebuild(ctx.model)

    return OpResult(success=True, message=f"Image '{ref}' removed", prefix="-")


HANDLERS: dict[str, callable] = {
    "image": op_image,
}
