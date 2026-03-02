"""Text content handlers — text, placeholder, bullet."""

from __future__ import annotations

from pptx.util import Pt

from fcp_core import OpResult, ParsedOp

from fcp_slides.server.resolvers import (
    SlidesOpContext,
    require_active_slide,
    resolve_shape_on_slide,
    resolve_slide,
)


# Placeholder type name → placeholder idx
_PLACEHOLDER_MAP = {
    "title": 0,
    "body": 1,
    "subtitle": 10,
    "center-title": 0,
    "slide-number": 12,
    "date": 13,
    "footer": 14,
}


def op_text(op: ParsedOp, ctx: SlidesOpContext) -> OpResult:
    """Set, append, or clear text content of a shape.

    Syntax:
      text set SHAPE TEXT
      text append SHAPE TEXT
      text clear SHAPE
    """
    if not op.positionals:
        return OpResult(success=False, message="Usage: text set|append|clear SHAPE [TEXT]")

    action = op.positionals[0].lower()
    rest = op.positionals[1:]

    # Determine target slide
    on_ref = op.params.get("on")
    if on_ref:
        result = resolve_slide(on_ref, ctx)
        if result is None:
            return OpResult(success=False, message=f"Slide not found: {on_ref!r}")
        slide, slide_idx = result
    else:
        active = require_active_slide(ctx)
        if isinstance(active, str):
            return OpResult(success=False, message=active)
        slide, slide_idx = active

    if action == "set":
        if len(rest) < 2:
            return OpResult(success=False, message='Usage: text set SHAPE "TEXT"')
        shape_ref = rest[0]
        text = rest[1]
        return _text_set(shape_ref, text, slide, slide_idx, ctx)

    elif action == "append":
        if len(rest) < 2:
            return OpResult(success=False, message='Usage: text append SHAPE "TEXT"')
        shape_ref = rest[0]
        text = rest[1]
        return _text_append(shape_ref, text, slide, slide_idx, ctx)

    elif action == "clear":
        if not rest:
            return OpResult(success=False, message="Usage: text clear SHAPE")
        shape_ref = rest[0]
        return _text_clear(shape_ref, slide, slide_idx, ctx)

    return OpResult(success=False, message=f"Unknown text action: {action!r}. Use: set, append, clear")


def _text_set(
    shape_ref: str, text: str, slide, slide_idx: int, ctx: SlidesOpContext
) -> OpResult:
    """Replace all text in a shape."""
    shape = resolve_shape_on_slide(shape_ref, slide, slide_idx, ctx)
    if shape is None:
        return OpResult(success=False, message=f"Shape not found: {shape_ref!r}")

    if not shape.has_text_frame:
        return OpResult(success=False, message=f"Shape '{shape_ref}' has no text frame")

    tf = shape.text_frame
    tf.clear()
    tf.text = text

    preview = text[:40] + ("..." if len(text) > 40 else "")
    return OpResult(success=True, message=f"Text set on '{shape_ref}': \"{preview}\"", prefix="*")


def _text_append(
    shape_ref: str, text: str, slide, slide_idx: int, ctx: SlidesOpContext
) -> OpResult:
    """Append text as a new paragraph."""
    shape = resolve_shape_on_slide(shape_ref, slide, slide_idx, ctx)
    if shape is None:
        return OpResult(success=False, message=f"Shape not found: {shape_ref!r}")

    if not shape.has_text_frame:
        return OpResult(success=False, message=f"Shape '{shape_ref}' has no text frame")

    tf = shape.text_frame
    p = tf.add_paragraph()
    p.text = text

    return OpResult(success=True, message=f"Text appended to '{shape_ref}'", prefix="*")


def _text_clear(
    shape_ref: str, slide, slide_idx: int, ctx: SlidesOpContext
) -> OpResult:
    """Clear all text from a shape."""
    shape = resolve_shape_on_slide(shape_ref, slide, slide_idx, ctx)
    if shape is None:
        return OpResult(success=False, message=f"Shape not found: {shape_ref!r}")

    if not shape.has_text_frame:
        return OpResult(success=False, message=f"Shape '{shape_ref}' has no text frame")

    shape.text_frame.clear()
    return OpResult(success=True, message=f"Text cleared on '{shape_ref}'", prefix="-")


def op_placeholder(op: ParsedOp, ctx: SlidesOpContext) -> OpResult:
    """Set text in a slide placeholder by type name.

    Syntax: placeholder set title|subtitle|body TEXT [on:SLIDE]
    """
    if not op.positionals:
        return OpResult(success=False, message="Usage: placeholder set title|subtitle|body TEXT")

    action = op.positionals[0].lower()
    if action != "set":
        return OpResult(success=False, message=f"Unknown placeholder action: {action!r}. Use: set")

    rest = op.positionals[1:]
    if len(rest) < 2:
        return OpResult(success=False, message='Usage: placeholder set TYPE "TEXT"')

    ph_type = rest[0].lower()
    text = rest[1]

    # Determine target slide
    on_ref = op.params.get("on")
    if on_ref:
        result = resolve_slide(on_ref, ctx)
        if result is None:
            return OpResult(success=False, message=f"Slide not found: {on_ref!r}")
        slide, slide_idx = result
    else:
        active = require_active_slide(ctx)
        if isinstance(active, str):
            return OpResult(success=False, message=active)
        slide, slide_idx = active

    ph_idx = _PLACEHOLDER_MAP.get(ph_type)
    if ph_idx is None:
        available = ", ".join(sorted(_PLACEHOLDER_MAP.keys()))
        return OpResult(success=False, message=f"Unknown placeholder type: {ph_type!r}. Use: {available}")

    # Find the placeholder on the slide
    ph_shape = None
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == ph_idx:
            ph_shape = shape
            break

    if ph_shape is None:
        # Try alternates: subtitle can also be idx 1 (body) on title slides
        alternates = {
            "title": (0, 15),
            "subtitle": (1, 10),
        }
        alt_indices = alternates.get(ph_type)
        if alt_indices:
            for shape in slide.placeholders:
                if shape.placeholder_format.idx in alt_indices:
                    ph_shape = shape
                    break

    if ph_shape is None:
        available = [
            f"{_get_ph_name(s.placeholder_format.idx)}(idx={s.placeholder_format.idx})"
            for s in slide.placeholders
        ]
        return OpResult(
            success=False,
            message=f"Placeholder '{ph_type}' not found on slide {slide_idx + 1}. Available: {', '.join(available)}",
        )

    ph_shape.text = text

    preview = text[:40] + ("..." if len(text) > 40 else "")
    return OpResult(
        success=True,
        message=f"Placeholder '{ph_type}' set: \"{preview}\"",
        prefix="*",
    )


def _get_ph_name(idx: int) -> str:
    """Get a friendly name for a placeholder index."""
    for name, i in _PLACEHOLDER_MAP.items():
        if i == idx:
            return name
    return f"ph{idx}"


def op_bullet(op: ParsedOp, ctx: SlidesOpContext) -> OpResult:
    """Add a bullet item to a shape at a given indent level.

    Syntax: bullet SHAPE TEXT [level:N]
    """
    if len(op.positionals) < 2:
        return OpResult(success=False, message='Usage: bullet SHAPE "TEXT" [level:N]')

    shape_ref = op.positionals[0]
    text = op.positionals[1]
    level = 0
    if "level" in op.params:
        try:
            level = int(op.params["level"])
        except ValueError:
            return OpResult(success=False, message=f"Invalid level: {op.params['level']!r}")

    # Determine target slide
    on_ref = op.params.get("on")
    if on_ref:
        result = resolve_slide(on_ref, ctx)
        if result is None:
            return OpResult(success=False, message=f"Slide not found: {on_ref!r}")
        slide, slide_idx = result
    else:
        active = require_active_slide(ctx)
        if isinstance(active, str):
            return OpResult(success=False, message=active)
        slide, slide_idx = active

    shape = resolve_shape_on_slide(shape_ref, slide, slide_idx, ctx)
    if shape is None:
        return OpResult(success=False, message=f"Shape not found: {shape_ref!r}")

    if not shape.has_text_frame:
        return OpResult(success=False, message=f"Shape '{shape_ref}' has no text frame")

    tf = shape.text_frame

    # If the text frame has only one empty paragraph, use it
    if len(tf.paragraphs) == 1 and not tf.paragraphs[0].text.strip():
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()

    p.text = text
    p.level = level

    return OpResult(
        success=True,
        message=f"Bullet added to '{shape_ref}' L{level}: \"{text[:40]}\"",
        prefix="+",
    )


HANDLERS: dict[str, callable] = {
    "text": op_text,
    "placeholder": op_placeholder,
    "bullet": op_bullet,
}
