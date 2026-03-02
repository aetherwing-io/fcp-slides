"""Shape management handlers — add, remove, move, resize, duplicate, textbox, connector."""

from __future__ import annotations

from copy import deepcopy

from pptx.util import Inches, Emu

from fcp_core import OpResult, ParsedOp

from fcp_slides.lib.shape_types import resolve_shape_type, list_shape_types
from fcp_slides.lib.units import parse_length
from fcp_slides.model.index import _shape_type_name
from fcp_slides.model.refs import ShapeRef
from fcp_slides.server.resolvers import (
    SlidesOpContext,
    extract_position,
    require_active_slide,
    resolve_shape_on_slide,
)


# Default dimensions for new shapes
_DEFAULT_LEFT = Inches(2)
_DEFAULT_TOP = Inches(2)
_DEFAULT_WIDTH = Inches(4)
_DEFAULT_HEIGHT = Inches(3)
_DEFAULT_TEXTBOX_WIDTH = Inches(6)
_DEFAULT_TEXTBOX_HEIGHT = Inches(1)


def op_shape(op: ParsedOp, ctx: SlidesOpContext) -> OpResult:
    """Dispatch shape sub-commands."""
    if not op.positionals:
        return OpResult(success=False, message="Usage: shape add|remove|move|resize|duplicate ...")

    action = op.positionals[0].lower()
    rest = op.positionals[1:]

    dispatch = {
        "add": _shape_add,
        "remove": _shape_remove,
        "move": _shape_move,
        "resize": _shape_resize,
        "duplicate": _shape_duplicate,
    }

    handler = dispatch.get(action)
    if handler is None:
        return OpResult(
            success=False,
            message=f"Unknown shape action: {action!r}. Use: add, remove, move, resize, duplicate",
        )

    return handler(rest, op.params, ctx)


def _shape_add(
    args: list[str], params: dict[str, str], ctx: SlidesOpContext
) -> OpResult:
    """Add a shape to the active slide."""
    result = require_active_slide(ctx)
    if isinstance(result, str):
        return OpResult(success=False, message=result)
    slide, slide_idx = result

    if not args:
        return OpResult(
            success=False,
            message=f"Usage: shape add TYPE [label:NAME] [x:POS] [y:POS] [w:SIZE] [h:SIZE]\nTypes: {', '.join(list_shape_types()[:10])}...",
        )

    shape_type_name = args[0]
    shape_type = resolve_shape_type(shape_type_name)
    if shape_type is None:
        return OpResult(
            success=False,
            message=f"Unknown shape type: {shape_type_name!r}. Use: {', '.join(list_shape_types()[:10])}...",
        )

    pos = extract_position(params)
    left = pos.get("left", _DEFAULT_LEFT)
    top = pos.get("top", _DEFAULT_TOP)
    width = pos.get("width", _DEFAULT_WIDTH)
    height = pos.get("height", _DEFAULT_HEIGHT)

    shape = slide.shapes.add_shape(shape_type, left, top, width, height)

    label = params.get("label", "")
    if label:
        ref = ShapeRef(
            label=label,
            slide_idx=slide_idx,
            shape_id=shape.shape_id,
            shape_type=shape_type_name,
        )
        ctx.index.add_shape_label(label, ref)
    else:
        ctx.index.rebuild(ctx.model)

    return OpResult(
        success=True,
        message=f"Shape '{shape_type_name}' added on slide {slide_idx + 1}" + (f" label:{label}" if label else ""),
        prefix="+",
    )


def _shape_remove(
    args: list[str], params: dict[str, str], ctx: SlidesOpContext
) -> OpResult:
    """Remove a shape from the active slide."""
    result = require_active_slide(ctx)
    if isinstance(result, str):
        return OpResult(success=False, message=result)
    slide, slide_idx = result

    if not args:
        return OpResult(success=False, message="Usage: shape remove SHAPE_REF")

    ref = args[0]
    shape = resolve_shape_on_slide(ref, slide, slide_idx, ctx)
    if shape is None:
        return OpResult(success=False, message=f"Shape not found: {ref!r}")

    # Remove shape via XML
    sp = shape._element
    sp.getparent().remove(sp)

    ctx.index.remove_shape_label(ref)
    ctx.index.rebuild(ctx.model)

    return OpResult(success=True, message=f"Shape '{ref}' removed", prefix="-")


def _shape_move(
    args: list[str], params: dict[str, str], ctx: SlidesOpContext
) -> OpResult:
    """Move a shape to a new position."""
    result = require_active_slide(ctx)
    if isinstance(result, str):
        return OpResult(success=False, message=result)
    slide, slide_idx = result

    if not args:
        return OpResult(success=False, message="Usage: shape move SHAPE_REF x:POS y:POS")

    ref = args[0]
    shape = resolve_shape_on_slide(ref, slide, slide_idx, ctx)
    if shape is None:
        return OpResult(success=False, message=f"Shape not found: {ref!r}")

    pos = extract_position(params)
    if "left" in pos:
        shape.left = pos["left"]
    if "top" in pos:
        shape.top = pos["top"]

    return OpResult(success=True, message=f"Shape '{ref}' moved", prefix="*")


def _shape_resize(
    args: list[str], params: dict[str, str], ctx: SlidesOpContext
) -> OpResult:
    """Resize a shape."""
    result = require_active_slide(ctx)
    if isinstance(result, str):
        return OpResult(success=False, message=result)
    slide, slide_idx = result

    if not args:
        return OpResult(success=False, message="Usage: shape resize SHAPE_REF w:SIZE h:SIZE")

    ref = args[0]
    shape = resolve_shape_on_slide(ref, slide, slide_idx, ctx)
    if shape is None:
        return OpResult(success=False, message=f"Shape not found: {ref!r}")

    pos = extract_position(params)
    if "width" in pos:
        shape.width = pos["width"]
    if "height" in pos:
        shape.height = pos["height"]
    if "left" in pos:
        shape.left = pos["left"]
    if "top" in pos:
        shape.top = pos["top"]

    return OpResult(success=True, message=f"Shape '{ref}' resized", prefix="*")


def _shape_duplicate(
    args: list[str], params: dict[str, str], ctx: SlidesOpContext
) -> OpResult:
    """Duplicate a shape on the same slide."""
    result = require_active_slide(ctx)
    if isinstance(result, str):
        return OpResult(success=False, message=result)
    slide, slide_idx = result

    if not args:
        return OpResult(success=False, message="Usage: shape duplicate SHAPE_REF [label:NAME]")

    ref = args[0]
    shape = resolve_shape_on_slide(ref, slide, slide_idx, ctx)
    if shape is None:
        return OpResult(success=False, message=f"Shape not found: {ref!r}")

    # Deep copy the XML element
    new_elem = deepcopy(shape._element)
    slide.shapes._spTree.append(new_elem)

    # Offset the copy slightly
    if shape.left is not None:
        new_elem_shapes = [s for s in slide.shapes if s._element is new_elem]
        if new_elem_shapes:
            new_shape = new_elem_shapes[0]
            new_shape.left = shape.left + Inches(0.5)
            new_shape.top = shape.top + Inches(0.5)

    label = params.get("label", "")
    ctx.index.rebuild(ctx.model)

    return OpResult(
        success=True,
        message=f"Shape '{ref}' duplicated" + (f" label:{label}" if label else ""),
        prefix="+",
    )


def op_textbox(op: ParsedOp, ctx: SlidesOpContext) -> OpResult:
    """Add a textbox with content (convenience verb).

    Syntax: textbox TEXT [x:POS] [y:POS] [w:SIZE] [h:SIZE] [label:NAME]
    """
    result = require_active_slide(ctx)
    if isinstance(result, str):
        return OpResult(success=False, message=result)
    slide, slide_idx = result

    if not op.positionals:
        return OpResult(success=False, message='Usage: textbox "TEXT" [x:POS] [y:POS] [w:SIZE] [h:SIZE]')

    text = op.positionals[0]

    pos = extract_position(op.params)
    left = pos.get("left", _DEFAULT_LEFT)
    top = pos.get("top", _DEFAULT_TOP)
    width = pos.get("width", _DEFAULT_TEXTBOX_WIDTH)
    height = pos.get("height", _DEFAULT_TEXTBOX_HEIGHT)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = text

    # Apply optional text styling
    if any(k in op.params for k in ("font", "size", "color", "bold", "align")):
        from pptx.util import Pt
        from fcp_slides.lib.colors import to_rgb
        for para in tf.paragraphs:
            for run in para.runs:
                if "font" in op.params:
                    run.font.name = op.params["font"]
                if "size" in op.params:
                    run.font.size = Pt(float(op.params["size"]))
                if "color" in op.params:
                    run.font.color.rgb = to_rgb(op.params["color"])
                if "bold" in op.params:
                    run.font.bold = True
            if "align" in op.params:
                from pptx.enum.text import PP_ALIGN
                align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
                align = align_map.get(op.params["align"].lower())
                if align:
                    para.alignment = align

    label = op.params.get("label", "")
    if label:
        ref = ShapeRef(
            label=label,
            slide_idx=slide_idx,
            shape_id=txBox.shape_id,
            shape_type="text_box",
        )
        ctx.index.add_shape_label(label, ref)
    else:
        ctx.index.rebuild(ctx.model)

    preview = text[:30] + ("..." if len(text) > 30 else "")
    return OpResult(
        success=True,
        message=f"Textbox added: \"{preview}\"" + (f" label:{label}" if label else ""),
        prefix="+",
    )


def op_connector(op: ParsedOp, ctx: SlidesOpContext) -> OpResult:
    """Connect two shapes with a connector line.

    Syntax: connector FROM_SHAPE TO_SHAPE [type:straight|elbow|curved]
    """
    result = require_active_slide(ctx)
    if isinstance(result, str):
        return OpResult(success=False, message=result)
    slide, slide_idx = result

    if len(op.positionals) < 2:
        return OpResult(success=False, message="Usage: connector FROM_SHAPE TO_SHAPE [type:straight|elbow|curved]")

    from_ref = op.positionals[0]
    to_ref = op.positionals[1]

    from_shape = resolve_shape_on_slide(from_ref, slide, slide_idx, ctx)
    if from_shape is None:
        return OpResult(success=False, message=f"Source shape not found: {from_ref!r}")

    to_shape = resolve_shape_on_slide(to_ref, slide, slide_idx, ctx)
    if to_shape is None:
        return OpResult(success=False, message=f"Target shape not found: {to_ref!r}")

    # Calculate connector endpoints from shape centers
    from_cx = (from_shape.left or 0) + (from_shape.width or 0) // 2
    from_cy = (from_shape.top or 0) + (from_shape.height or 0) // 2
    to_cx = (to_shape.left or 0) + (to_shape.width or 0) // 2
    to_cy = (to_shape.top or 0) + (to_shape.height or 0) // 2

    # Add as a line (python-pptx doesn't have native connector API)
    from pptx.util import Emu
    connector = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR_TYPE.STRAIGHT
        from_cx, from_cy,
        to_cx, to_cy,
    )

    return OpResult(
        success=True,
        message=f"Connected '{from_ref}' → '{to_ref}'",
        prefix="~",
    )


HANDLERS: dict[str, callable] = {
    "shape": op_shape,
    "textbox": op_textbox,
    "connector": op_connector,
}
