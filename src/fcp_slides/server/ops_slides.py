"""Slide management handlers — add, remove, rename, copy, move, hide, unhide, activate."""

from __future__ import annotations

from copy import deepcopy
from lxml import etree

from pptx.util import Emu

from fcp_core import OpResult, ParsedOp

from fcp_slides.lib.layout_names import resolve_layout, list_layouts
from fcp_slides.server.resolvers import SlidesOpContext, resolve_slide


def op_slide(op: ParsedOp, ctx: SlidesOpContext) -> OpResult:
    """Dispatch slide sub-commands."""
    if not op.positionals:
        return OpResult(success=False, message="Usage: slide add|remove|rename|copy|move|hide|unhide|activate")

    action = op.positionals[0].lower()
    rest = op.positionals[1:]

    dispatch = {
        "add": _slide_add,
        "remove": _slide_remove,
        "rename": _slide_rename,
        "copy": _slide_copy,
        "move": _slide_move,
        "hide": _slide_hide,
        "unhide": _slide_unhide,
        "activate": _slide_activate,
    }

    handler = dispatch.get(action)
    if handler is None:
        return OpResult(
            success=False,
            message=f"Unknown slide action: {action!r}. Use: add, remove, rename, copy, move, hide, unhide, activate",
        )

    return handler(rest, op.params, ctx)


def _slide_add(
    args: list[str], params: dict[str, str], ctx: SlidesOpContext
) -> OpResult:
    """Add a new slide with optional layout."""
    layout_name = params.get("layout", "blank")
    layout = resolve_layout(layout_name, ctx.prs)
    if layout is None:
        available = ", ".join(list_layouts(ctx.prs))
        return OpResult(
            success=False,
            message=f"Layout not found: {layout_name!r}. Available: {available}",
        )

    slide = ctx.prs.slides.add_slide(layout)
    slide_idx = len(ctx.prs.slides) - 1

    # Auto-label
    label = params.get("label", f"s{slide_idx + 1}")
    ctx.index.add_slide_label(label, slide_idx)
    ctx.index.active_slide = slide_idx

    # Rebuild shape index for the new slide
    ctx.index.rebuild(ctx.model)

    return OpResult(
        success=True,
        message=f"Slide {slide_idx + 1} added [{layout.name}] label:{label}",
        prefix="+",
    )


def _slide_remove(
    args: list[str], params: dict[str, str], ctx: SlidesOpContext
) -> OpResult:
    """Remove a slide."""
    if not args:
        return OpResult(success=False, message="Usage: slide remove SLIDE_REF")

    ref = args[0]
    result = resolve_slide(ref, ctx)
    if result is None:
        return OpResult(success=False, message=f"Slide not found: {ref!r}")

    slide, slide_idx = result

    if len(ctx.prs.slides) <= 1:
        return OpResult(success=False, message="Cannot remove the last slide")

    # Remove slide via XML manipulation (not in public API)
    rId = ctx.prs.slides._sldIdLst[slide_idx].get(
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    )
    # Remove from slide ID list
    sldIdLst = ctx.prs.slides._sldIdLst
    sldId = sldIdLst[slide_idx]
    sldIdLst.remove(sldId)

    # Remove the relationship
    ctx.prs.part.drop_rel(rId)

    # Update index
    was_active = ctx.index.active_slide == slide_idx
    ctx.index.rebuild(ctx.model)
    if was_active:
        ctx.index.active_slide = min(slide_idx, len(ctx.prs.slides) - 1)

    return OpResult(success=True, message=f"Slide {slide_idx + 1} removed", prefix="-")


def _slide_rename(
    args: list[str], params: dict[str, str], ctx: SlidesOpContext
) -> OpResult:
    """Rename a slide label."""
    if len(args) < 2:
        return OpResult(success=False, message="Usage: slide rename SLIDE_REF NEW_LABEL")

    ref = args[0]
    new_label = args[1]

    slide_idx = ctx.index.resolve_slide_idx(ref, ctx.model)
    if slide_idx is None:
        return OpResult(success=False, message=f"Slide not found: {ref!r}")

    # Remove old label and add new one
    ctx.index.remove_slide_label(ref)
    ctx.index.add_slide_label(new_label, slide_idx)

    return OpResult(
        success=True,
        message=f"Slide {slide_idx + 1} label: {new_label}",
        prefix="*",
    )


def _slide_copy(
    args: list[str], params: dict[str, str], ctx: SlidesOpContext
) -> OpResult:
    """Copy a slide (deep copy via XML)."""
    if not args:
        return OpResult(success=False, message="Usage: slide copy SLIDE_REF [label:NAME]")

    ref = args[0]
    result = resolve_slide(ref, ctx)
    if result is None:
        return OpResult(success=False, message=f"Slide not found: {ref!r}")

    slide, slide_idx = result

    # Deep copy: duplicate the slide's XML and relationships
    layout = slide.slide_layout
    new_slide = ctx.prs.slides.add_slide(layout)
    new_idx = len(ctx.prs.slides) - 1

    # Copy shapes by duplicating the slide's XML body
    for shape_elem in slide.shapes._spTree:
        new_slide.shapes._spTree.append(deepcopy(shape_elem))

    label = params.get("label", f"s{new_idx + 1}")
    ctx.index.add_slide_label(label, new_idx)
    ctx.index.active_slide = new_idx
    ctx.index.rebuild(ctx.model)

    return OpResult(
        success=True,
        message=f"Slide {slide_idx + 1} copied as slide {new_idx + 1} label:{label}",
        prefix="+",
    )


def _slide_move(
    args: list[str], params: dict[str, str], ctx: SlidesOpContext
) -> OpResult:
    """Move a slide to a new position."""
    if not args:
        return OpResult(success=False, message="Usage: slide move SLIDE_REF to:N | after:REF | before:REF")

    ref = args[0]
    slide_idx = ctx.index.resolve_slide_idx(ref, ctx.model)
    if slide_idx is None:
        return OpResult(success=False, message=f"Slide not found: {ref!r}")

    # Determine target position
    target = None
    if "to" in params:
        try:
            target = int(params["to"]) - 1  # Convert to 0-based
        except ValueError:
            return OpResult(success=False, message=f"Invalid position: {params['to']!r}")
    elif "after" in params:
        after_idx = ctx.index.resolve_slide_idx(params["after"], ctx.model)
        if after_idx is None:
            return OpResult(success=False, message=f"Slide not found: {params['after']!r}")
        target = after_idx + 1
    elif "before" in params:
        before_idx = ctx.index.resolve_slide_idx(params["before"], ctx.model)
        if before_idx is None:
            return OpResult(success=False, message=f"Slide not found: {params['before']!r}")
        target = before_idx

    if target is None:
        return OpResult(success=False, message="Specify position with to:N, after:REF, or before:REF")

    slide_count = len(ctx.prs.slides)
    target = max(0, min(target, slide_count - 1))

    if target == slide_idx:
        return OpResult(success=True, message=f"Slide {slide_idx + 1} already at position {target + 1}", prefix="*")

    # Move via XML manipulation
    sldIdLst = ctx.prs.slides._sldIdLst
    elem = sldIdLst[slide_idx]
    sldIdLst.remove(elem)
    sldIdLst.insert(target, elem)

    ctx.index.rebuild(ctx.model)
    ctx.index.active_slide = target

    return OpResult(
        success=True,
        message=f"Slide moved from {slide_idx + 1} to {target + 1}",
        prefix="*",
    )


def _slide_hide(
    args: list[str], params: dict[str, str], ctx: SlidesOpContext
) -> OpResult:
    """Hide a slide (XML attribute manipulation)."""
    if not args:
        return OpResult(success=False, message="Usage: slide hide SLIDE_REF")

    ref = args[0]
    result = resolve_slide(ref, ctx)
    if result is None:
        return OpResult(success=False, message=f"Slide not found: {ref!r}")

    slide, slide_idx = result

    # Set show attribute to 0 in the slide XML
    slide_elem = slide._element
    slide_elem.set("show", "0")

    return OpResult(success=True, message=f"Slide {slide_idx + 1} hidden", prefix="*")


def _slide_unhide(
    args: list[str], params: dict[str, str], ctx: SlidesOpContext
) -> OpResult:
    """Unhide a slide."""
    if not args:
        return OpResult(success=False, message="Usage: slide unhide SLIDE_REF")

    ref = args[0]
    result = resolve_slide(ref, ctx)
    if result is None:
        return OpResult(success=False, message=f"Slide not found: {ref!r}")

    slide, slide_idx = result
    slide_elem = slide._element
    if "show" in slide_elem.attrib:
        del slide_elem.attrib["show"]

    return OpResult(success=True, message=f"Slide {slide_idx + 1} unhidden", prefix="*")


def _slide_activate(
    args: list[str], params: dict[str, str], ctx: SlidesOpContext
) -> OpResult:
    """Switch active slide."""
    if not args:
        return OpResult(success=False, message="Usage: slide activate SLIDE_REF")

    ref = args[0]
    slide_idx = ctx.index.resolve_slide_idx(ref, ctx.model)
    if slide_idx is None:
        return OpResult(success=False, message=f"Slide not found: {ref!r}")

    ctx.index.active_slide = slide_idx
    return OpResult(success=True, message=f"Active slide: {slide_idx + 1}", prefix="*")


HANDLERS: dict[str, callable] = {
    "slide": op_slide,
}
