"""Byte-snapshot undo/redo for python-pptx presentations.

SnapshotEvent captures before/after states as bytes from
prs.save(BytesIO) / Presentation(BytesIO). This enables
simple, correct undo/redo for all presentation mutations.
"""

from __future__ import annotations

from dataclasses import dataclass, field
from io import BytesIO

from pptx import Presentation


@dataclass
class SnapshotEvent:
    """Event type for byte-snapshot undo/redo."""

    type: str = "snapshot"
    before: bytes = field(default=b"", repr=False)
    after: bytes = field(default=b"", repr=False)
    summary: str = ""


def snapshot_presentation(prs: Presentation) -> bytes:
    """Serialize a presentation to bytes via BytesIO."""
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def restore_presentation(data: bytes) -> Presentation:
    """Deserialize a presentation from snapshot bytes."""
    buf = BytesIO(data)
    return Presentation(buf)


class SlidesModel:
    """Thin wrapper around python-pptx Presentation for in-place undo/redo.

    The session dispatcher holds a reference to this object.
    reverse_event/replay_event replace self.prs in place so
    the session reference stays valid.
    """

    def __init__(self, title: str = "Untitled", prs: Presentation | None = None):
        self.title = title
        self.prs: Presentation = prs or Presentation()
        self.file_path: str | None = None

    def snapshot(self) -> bytes:
        """Take a byte snapshot of the current presentation state."""
        return snapshot_presentation(self.prs)

    def restore(self, data: bytes) -> None:
        """Replace the presentation from snapshot bytes (in-place for undo/redo)."""
        self.prs = restore_presentation(data)
