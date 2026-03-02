"""SlideIndex — lightweight metadata index for slide/shape resolution.

Tracks slide labels, shape labels, and active slide context.
Rebuild is O(slides * shapes), not O(content).
"""

from __future__ import annotations

from typing import TYPE_CHECKING

from fcp_slides.model.refs import ShapeRef

if TYPE_CHECKING:
    from pptx.shapes.base import BaseShape
    from pptx.slide import Slide

    from fcp_slides.model.snapshot import SlidesModel


# Map python-pptx MSO_SHAPE_TYPE enum values to friendly names
_SHAPE_TYPE_NAMES = {
    1: "auto_shape",
    5: "freeform",
    6: "group",
    13: "picture",
    14: "placeholder",
    17: "text_box",
    19: "table",
    3: "chart",
}

# Map placeholder idx to common names
_PLACEHOLDER_NAMES = {
    0: "title",
    1: "body",
    10: "subtitle",
    12: "slide_number",
    13: "date",
    14: "footer",
}


def _shape_type_name(shape: BaseShape) -> str:
    """Get a friendly type name for a shape."""
    from pptx.util import Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    st = shape.shape_type
    if st is not None:
        val = st if isinstance(st, int) else st.real
        if val in _SHAPE_TYPE_NAMES:
            return _SHAPE_TYPE_NAMES[val]
    # Fallback: check for specific types
    if shape.has_table:
        return "table"
    if shape.has_chart:
        return "chart"
    if shape.has_text_frame:
        return "text_box"
    return "shape"


def _auto_label_shape(shape: BaseShape, slide_idx: int, shape_counter: dict[str, int]) -> str:
    """Generate an automatic label for a shape."""
    if shape.has_text_frame and shape.text_frame.text.strip():
        # Use first 20 chars of text content as hint
        text = shape.text_frame.text.strip()[:20].lower()
        # Sanitize: keep alphanumeric and spaces, replace spaces with _
        sanitized = "".join(c if c.isalnum() or c == " " else "" for c in text)
        sanitized = sanitized.strip().replace(" ", "_")
        if sanitized:
            return sanitized

    type_name = _shape_type_name(shape)
    key = f"s{slide_idx + 1}_{type_name}"
    count = shape_counter.get(key, 0) + 1
    shape_counter[key] = count
    return f"{key}{count}"


def _placeholder_label(shape: BaseShape, slide_idx: int) -> str | None:
    """Generate a label for a placeholder shape."""
    if not shape.is_placeholder:
        return None
    ph = shape.placeholder_format
    if ph is None:
        return None
    idx = ph.idx
    name = _PLACEHOLDER_NAMES.get(idx)
    if name:
        return f"s{slide_idx + 1}_{name}"
    return f"s{slide_idx + 1}_ph{idx}"


class SlideIndex:
    """Index for O(1) slide and shape resolution by label."""

    def __init__(self) -> None:
        self._slide_labels: dict[str, int] = {}  # label → 0-based slide index
        self._shape_labels: dict[str, ShapeRef] = {}  # label → ShapeRef
        self._active_slide: int = 0  # 0-based index

    @property
    def active_slide(self) -> int:
        return self._active_slide

    @active_slide.setter
    def active_slide(self, idx: int) -> None:
        self._active_slide = idx

    def slide_count(self, model: SlidesModel) -> int:
        """Return the number of slides in the presentation."""
        return len(model.prs.slides)

    def resolve_slide_idx(self, ref: str, model: SlidesModel) -> int | None:
        """Resolve a slide reference to a 0-based index.

        Accepts:
          - Label: "overview", "s1"
          - 1-based number: "1", "3"
          - Keyword: "active", "last"
        """
        ref = ref.strip()
        low = ref.lower()
        slide_count = len(model.prs.slides)

        if low == "active":
            return self._active_slide if self._active_slide < slide_count else None

        if low == "last":
            return slide_count - 1 if slide_count > 0 else None

        # Check labels
        if low in self._slide_labels:
            idx = self._slide_labels[low]
            return idx if idx < slide_count else None

        # Try 1-based number
        try:
            n = int(ref)
            if 1 <= n <= slide_count:
                return n - 1
            return None
        except ValueError:
            pass

        return None

    def resolve_shape(self, ref: str, slide_idx: int | None = None) -> ShapeRef | None:
        """Resolve a shape reference to a ShapeRef.

        Accepts labels. If slide_idx is given, constrains to that slide.
        """
        low = ref.strip().lower()
        shape_ref = self._shape_labels.get(low)
        if shape_ref is None:
            return None
        if slide_idx is not None and shape_ref.slide_idx != slide_idx:
            return None
        return shape_ref

    def get_shapes_on_slide(self, slide_idx: int) -> list[ShapeRef]:
        """Get all shape refs on a given slide."""
        return [sr for sr in self._shape_labels.values() if sr.slide_idx == slide_idx]

    def add_slide_label(self, label: str, idx: int) -> None:
        """Register a slide label."""
        self._slide_labels[label.lower()] = idx

    def remove_slide_label(self, label: str) -> None:
        """Remove a slide label."""
        self._slide_labels.pop(label.lower(), None)

    def add_shape_label(self, label: str, ref: ShapeRef) -> None:
        """Register a shape label."""
        self._shape_labels[label.lower()] = ref

    def remove_shape_label(self, label: str) -> None:
        """Remove a shape label."""
        self._shape_labels.pop(label.lower(), None)

    def rebuild(self, model: SlidesModel) -> None:
        """Full rebuild from presentation state (after undo/redo/open).

        Scans all slides and shapes to rebuild label indices.
        """
        self._slide_labels.clear()
        self._shape_labels.clear()

        prs = model.prs
        shape_counter: dict[str, int] = {}

        for slide_idx, slide in enumerate(prs.slides):
            # Auto-label slides as s1, s2, ...
            slide_label = f"s{slide_idx + 1}"
            self._slide_labels[slide_label] = slide_idx

            for shape in slide.shapes:
                # Placeholder label takes priority
                ph_label = _placeholder_label(shape, slide_idx)
                if ph_label:
                    ref = ShapeRef(
                        label=ph_label,
                        slide_idx=slide_idx,
                        shape_id=shape.shape_id,
                        shape_type=_shape_type_name(shape),
                        placeholder_idx=shape.placeholder_format.idx if shape.placeholder_format else None,
                    )
                    self._shape_labels[ph_label.lower()] = ref

                # Auto-label for non-placeholders (or as additional label)
                auto = _auto_label_shape(shape, slide_idx, shape_counter)
                if auto.lower() not in self._shape_labels:
                    ref = ShapeRef(
                        label=auto,
                        slide_idx=slide_idx,
                        shape_id=shape.shape_id,
                        shape_type=_shape_type_name(shape),
                        placeholder_idx=shape.placeholder_format.idx if shape.is_placeholder and shape.placeholder_format else None,
                    )
                    self._shape_labels[auto.lower()] = ref

        # Clamp active slide
        slide_count = len(prs.slides)
        if self._active_slide >= slide_count:
            self._active_slide = max(0, slide_count - 1)

    def clear(self) -> None:
        """Reset index state."""
        self._slide_labels.clear()
        self._shape_labels.clear()
        self._active_slide = 0

    def shift_slides_after(self, after_idx: int, delta: int) -> None:
        """Shift slide indices after a given position (for insert/remove).

        Updates both slide labels and shape refs that point to shifted slides.
        """
        # Update slide labels
        updated: dict[str, int] = {}
        for label, idx in self._slide_labels.items():
            if idx > after_idx:
                updated[label] = idx + delta
            elif idx == after_idx and delta < 0:
                # This slide was removed — skip
                continue
            else:
                updated[label] = idx
        self._slide_labels = updated

        # Update shape refs
        to_remove: list[str] = []
        for label, ref in self._shape_labels.items():
            if ref.slide_idx > after_idx:
                ref.slide_idx += delta
            elif ref.slide_idx == after_idx and delta < 0:
                to_remove.append(label)
        for label in to_remove:
            del self._shape_labels[label]
