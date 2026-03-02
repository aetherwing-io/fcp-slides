"""SlidesAdapter — FcpDomainAdapter implementation for python-pptx presentations.

Bridges fcp-core to python-pptx via SlidesModel (thin wrapper for in-place
undo/redo). Handles batch atomicity and snapshot-based undo/redo.
"""

from __future__ import annotations

from pptx import Presentation

from fcp_core import EventLog, OpResult, ParsedOp

from fcp_slides.lib.units import format_length
from fcp_slides.model.index import SlideIndex
from fcp_slides.model.snapshot import SlidesModel, SnapshotEvent
from fcp_slides.server.queries import dispatch_query
from fcp_slides.server.resolvers import SlidesOpContext

# Import all handler dicts
from fcp_slides.server.ops_slides import HANDLERS as SLIDES_HANDLERS
from fcp_slides.server.ops_shapes import HANDLERS as SHAPES_HANDLERS
from fcp_slides.server.ops_text import HANDLERS as TEXT_HANDLERS
from fcp_slides.server.ops_tables import HANDLERS as TABLES_HANDLERS
from fcp_slides.server.ops_charts import HANDLERS as CHARTS_HANDLERS
from fcp_slides.server.ops_images import HANDLERS as IMAGES_HANDLERS
from fcp_slides.server.ops_layout import HANDLERS as LAYOUT_HANDLERS
from fcp_slides.server.ops_style import HANDLERS as STYLE_HANDLERS
from fcp_slides.server.ops_notes import HANDLERS as NOTES_HANDLERS

# Max snapshot events in undo history
MAX_EVENTS = 15


class SlidesAdapter:
    """FcpDomainAdapter[SlidesModel, SnapshotEvent] for presentation operations."""

    def __init__(self) -> None:
        self.index = SlideIndex()

        # Merge all verb handlers
        self._handlers: dict[str, callable] = {}
        for h in (
            SLIDES_HANDLERS, SHAPES_HANDLERS, TEXT_HANDLERS,
            TABLES_HANDLERS, CHARTS_HANDLERS, IMAGES_HANDLERS,
            LAYOUT_HANDLERS, STYLE_HANDLERS, NOTES_HANDLERS,
        ):
            self._handlers.update(h)

    # -- FcpDomainAdapter protocol --

    def create_empty(self, title: str, params: dict[str, str]) -> SlidesModel:
        """Create a new empty presentation."""
        prs = Presentation()
        model = SlidesModel(title=title, prs=prs)
        self.index.clear()
        return model

    def serialize(self, model: SlidesModel, path: str) -> None:
        """Save presentation to file."""
        model.prs.save(path)
        model.file_path = path

    def deserialize(self, path: str) -> SlidesModel:
        """Load presentation from file."""
        prs = Presentation(path)
        # Extract title from path
        title = path.rsplit("/", 1)[-1]
        model = SlidesModel(title=title, prs=prs)
        model.file_path = path
        self.index.rebuild(model)
        return model

    def rebuild_indices(self, model: SlidesModel) -> None:
        """Rebuild index after undo/redo."""
        self.index.rebuild(model)

    def get_digest(self, model: SlidesModel) -> str:
        """Return a compact state fingerprint."""
        prs = model.prs
        slide_count = len(prs.slides)
        shape_count = sum(len(list(s.shapes)) for s in prs.slides)
        active = self.index.active_slide + 1
        size = ""
        if prs.slide_width and prs.slide_height:
            size = f", Size: {format_length(prs.slide_width)}x{format_length(prs.slide_height)}"
        return f"Active: slide {active}, Slides: {slide_count}, Shapes: {shape_count}{size}"

    def dispatch_op(
        self, op: ParsedOp, model: SlidesModel, log: EventLog
    ) -> OpResult:
        """Execute a parsed operation on the model."""
        handler = self._handlers.get(op.verb)
        if handler is None:
            from fcp_core import suggest
            s = suggest(op.verb, list(self._handlers.keys()))
            msg = f"Unknown verb: {op.verb!r}"
            if s:
                msg += f"\n  try: {s}"
            return OpResult(success=False, message=msg)

        # Take pre-op snapshot
        before = model.snapshot()

        # Build context
        ctx = SlidesOpContext(
            prs=model.prs,
            index=self.index,
            model=model,
        )

        # Dispatch
        try:
            result = handler(op, ctx)
        except NotImplementedError as exc:
            return OpResult(success=False, message=str(exc))
        except (ValueError, KeyError, TypeError, AttributeError) as exc:
            return OpResult(success=False, message=f"Error: {exc}")

        if not result.success:
            return result

        # Log snapshot for undo
        after = model.snapshot()
        log.append(SnapshotEvent(before=before, after=after, summary=op.raw))

        return result

    def take_snapshot(self, model: SlidesModel) -> bytes:
        """Return byte snapshot for batch rollback."""
        return model.snapshot()

    def restore_snapshot(self, model: SlidesModel, snapshot: bytes) -> None:
        """Restore model from snapshot and rebuild indices."""
        model.restore(snapshot)
        self.rebuild_indices(model)

    def dispatch_query(self, query: str, model: SlidesModel) -> str:
        """Execute a query against the model."""
        return dispatch_query(query, model, self.index)

    def reverse_event(self, event: SnapshotEvent, model: SlidesModel) -> None:
        """Undo — restore from before-snapshot."""
        model.restore(event.before)
        self.index.rebuild(model)

    def replay_event(self, event: SnapshotEvent, model: SlidesModel) -> None:
        """Redo — restore from after-snapshot."""
        model.restore(event.after)
        self.index.rebuild(model)
